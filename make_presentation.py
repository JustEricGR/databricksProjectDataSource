"""
Generate DataProcessing Pipeline Architecture presentation.
Run: python make_presentation.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

# ── Palette ───────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0F, 0x11, 0x17)   # near-black
SURFACE   = RGBColor(0x1A, 0x1D, 0x27)   # dark card
BRONZE_C  = RGBColor(0xCD, 0x7F, 0x32)
SILVER_C  = RGBColor(0xA0, 0xA8, 0xC0)
GOLD_C    = RGBColor(0xF0, 0xB4, 0x29)
AGENT_C   = RGBColor(0x6E, 0x7B, 0xF0)   # indigo
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
MUTED     = RGBColor(0x8B, 0x90, 0xA8)
ACCENT    = RGBColor(0x4A, 0xC9, 0xFF)   # cyan

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank = prs.slide_layouts[6]  # completely blank

# ── Helper functions ──────────────────────────────────────────────────────────

def bg(slide, color=BG):
    bg_shape = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = color
    bg_shape.line.fill.background()

def box(slide, x, y, w, h, fill=SURFACE, alpha=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, size=20, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def divider(slide, y, color=MUTED, w=13.33):
    line = slide.shapes.add_shape(1, Inches(0), Inches(y), Inches(w), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

def badge(slide, label, x, y, color):
    b = box(slide, x, y, 1.5, 0.35, fill=color)
    txt(slide, label, x + 0.05, y + 0.03, 1.4, 0.3, size=11, bold=True,
        color=BG, align=PP_ALIGN.CENTER)

def arrow(slide, x, y, vertical=False):
    if vertical:
        arr = slide.shapes.add_shape(13, Inches(x), Inches(y), Inches(0.3), Inches(0.4))
    else:
        arr = slide.shapes.add_shape(13, Inches(x), Inches(y), Inches(0.5), Inches(0.3))
    arr.fill.solid()
    arr.fill.fore_color.rgb = MUTED
    arr.line.fill.background()

# ── Slide 1 — Title ───────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank)
bg(s1)
# Accent bar left
b = s1.shapes.add_shape(1, 0, 0, Inches(0.08), SLIDE_H)
b.fill.solid(); b.fill.fore_color.rgb = AGENT_C; b.line.fill.background()

txt(s1, "DataProcessing Pipeline", 0.4, 1.8, 12, 1.2, size=46, bold=True, color=WHITE)
txt(s1, "Architecture & Data Flow", 0.4, 3.0, 12, 0.8, size=30, color=SILVER_C)
divider(s1, 4.0, color=AGENT_C, w=12)
txt(s1, "Databricks DLT  ·  Unity Catalog  ·  AI Agent  ·  Lakeview Dashboards",
    0.4, 4.15, 12, 0.5, size=14, color=MUTED)
txt(s1, "JustEricGR / databricksProjectDataSource", 0.4, 6.8, 12, 0.4,
    size=11, color=MUTED, italic=True)

# ── Slide 2 — End-to-End Pipeline Overview ────────────────────────────────────
s2 = prs.slides.add_slide(blank)
bg(s2)
txt(s2, "End-to-End Pipeline Overview", 0.4, 0.3, 12, 0.6, size=28, bold=True, color=WHITE)
divider(s2, 1.05)

# Flow boxes
NODES = [
    ("GitHub\nCSV Files", 0.3,  2.6, MUTED),
    ("Bronze\nDLT Pipeline", 2.3, 2.6, BRONZE_C),
    ("Silver\nDLT Pipeline", 4.3, 2.6, SILVER_C),
    ("AI Agent\n(Llama 3.3)", 6.3, 2.6, AGENT_C),
    ("Gold\nDLT Views",  8.3,  2.6, GOLD_C),
    ("Lakeview\nDashboards", 10.3, 2.6, ACCENT),
]
for label, x, y, c in NODES:
    box(s2, x, y, 1.8, 1.1, fill=c)
    txt(s2, label, x+0.05, y+0.15, 1.7, 0.8, size=12, bold=True,
        color=BG, align=PP_ALIGN.CENTER)

# Arrows between nodes
for xi in [2.1, 4.1, 6.1, 8.1, 10.1]:
    txt(s2, "→", xi, 2.95, 0.25, 0.4, size=20, bold=True, color=MUTED,
        align=PP_ALIGN.CENTER)

# Job task bar
box(s2, 1.9, 4.1, 8.8, 0.5, fill=RGBColor(0x1E, 0x22, 0x35))
txt(s2, "Databricks Job:   bronze  →  silver  →  ai_agent  (agent triggers gold selectively)",
    2.0, 4.17, 8.6, 0.38, size=12, color=MUTED)

# Sub-labels
for label, x in [("dataSource/", 0.3), ("DLT incremental", 2.3),
                  ("DLT incremental", 4.3), ("selective refresh", 6.3),
                  ("75 views", 8.3), ("5 domains", 10.3)]:
    txt(s2, label, x, 3.8, 1.8, 0.3, size=10, color=MUTED, align=PP_ALIGN.CENTER)

# Key stats row
stats = [
    ("17 CSVs", "Source files"),
    ("17 tables", "Bronze"),
    ("17 tables", "Silver"),
    ("<5 min", "New table run"),
    ("75 views", "Gold analytics"),
    ("5 dashboards", "Lakeview"),
]
for i, (val, lbl) in enumerate(stats):
    bx = 0.3 + i * 2.1
    box(s2, bx, 5.2, 1.9, 0.8, fill=SURFACE)
    txt(s2, val, bx, 5.27, 1.9, 0.38, size=18, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s2, lbl, bx, 5.62, 1.9, 0.28, size=10, color=MUTED, align=PP_ALIGN.CENTER)

txt(s2, "500-table metastore limit enforced — quota checked before every operation",
    0.4, 6.5, 12.5, 0.4, size=11, color=GOLD_C, italic=True)

# ── Slide 3 — Data Ingestion (Bronze) ─────────────────────────────────────────
s3 = prs.slides.add_slide(blank)
bg(s3)
box(s3, 0, 0, 0.08, 7.5, fill=BRONZE_C).line.fill.background()
txt(s3, "Data Ingestion — Bronze Layer", 0.4, 0.3, 12, 0.6, size=28, bold=True, color=WHITE)
txt(s3, "Pipeline: data_ingestion_git  ·  DLT  ·  Unity Catalog: dataingestionproject.bronze",
    0.4, 0.9, 12.5, 0.35, size=12, color=MUTED)
divider(s3, 1.32)

# Left column — how it works
box(s3, 0.3, 1.5, 5.8, 4.8, fill=SURFACE)
txt(s3, "How It Works", 0.55, 1.6, 5.3, 0.4, size=15, bold=True, color=BRONZE_C)
steps = [
    ("1", "GitHub Actions detects pushed CSV(s) in dataSource/"),
    ("2", "Passes changed_files param to Databricks job"),
    ("3", "ingest.py fetches file list from GitHub Contents API"),
    ("4", "Incremental: only processes new/changed CSVs"),
    ("5", "Full fallback: re-ingests all CSVs on manual runs"),
    ("6", "pandas.read_csv (latin-1) → drop Unnamed: columns"),
    ("7", "Strip UTF-8 BOM from column names (ï»¿ → clean)"),
    ("8", "@dp.table creates a DLT managed table per CSV"),
]
for i, (num, step) in enumerate(steps):
    y = 2.1 + i * 0.47
    box(s3, 0.45, y, 0.3, 0.3, fill=BRONZE_C)
    txt(s3, num, 0.45, y, 0.3, 0.3, size=10, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(s3, step, 0.9, y, 4.9, 0.35, size=11, color=WHITE)

# Right column — tables & safeguards
box(s3, 6.5, 1.5, 6.5, 2.2, fill=SURFACE)
txt(s3, "Ingested Tables (17)", 6.75, 1.6, 6.0, 0.4, size=15, bold=True, color=BRONZE_C)
tables = ["dim_customers", "dim_products", "fact_sales", "cust_info", "cust_az12",
          "loc_a101", "prd_info", "px_cat_g1v2", "sales_details",
          "annual_enterprise_survey…(×2)", "business_operations…(×3)", "meets (new)"]
for i, tbl in enumerate(tables):
    col = i % 2; row = i // 2
    txt(s3, f"• {tbl}", 6.55 + col*3.2, 2.1 + row*0.38, 3.1, 0.35, size=10, color=SILVER_C)

box(s3, 6.5, 3.85, 6.5, 2.4, fill=SURFACE)
txt(s3, "Safeguards", 6.75, 3.95, 6.0, 0.4, size=15, bold=True, color=BRONZE_C)
guards = [
    ("Metastore quota check", "Stops if >= 490 tables in metastore"),
    ("Multi-file push support", "Space-separated changed_files handled"),
    ("BOM stripping",          "ï»¿ prefix removed from column names"),
    ("Unnamed: columns",       "Pandas trailing empty columns dropped"),
    ("Rate limit handling",     "GitHub API 403 → RuntimeError raised"),
]
for i, (k, v) in enumerate(guards):
    txt(s3, f"✓  {k}", 6.65, 4.45 + i*0.36, 2.8, 0.33, size=10, bold=True, color=GOLD_C)
    txt(s3, v, 9.45, 4.45 + i*0.36, 3.3, 0.33, size=10, color=MUTED)

# ── Slide 4 — Silver Transformations ──────────────────────────────────────────
s4 = prs.slides.add_slide(blank)
bg(s4)
box(s4, 0, 0, 0.08, 7.5, fill=SILVER_C).line.fill.background()
txt(s4, "Silver Transformations", 0.4, 0.3, 12, 0.6, size=28, bold=True, color=WHITE)
txt(s4, "Pipeline: silver_transformations  ·  DLT Materialized Views  ·  dataingestionproject.silver",
    0.4, 0.9, 12.5, 0.35, size=12, color=MUTED)
divider(s4, 1.32)

# Transformation rules
box(s4, 0.3, 1.5, 8.3, 4.9, fill=SURFACE)
txt(s4, "Transformation Rules Applied to Every Table", 0.55, 1.6, 8.0, 0.4,
    size=15, bold=True, color=SILVER_C)

rules = [
    (GOLD_C,    "Date Standardization",
     "standardize_date() — coalesces yyyy-MM-dd / MM/dd/yyyy / dd/MM/yyyy → dd/MM/yyyy string"),
    (ACCENT,    "Gender Normalization",
     "standardize_gender() — F/Female → 'Female', M/Male → 'Male', else NULL (rows dropped)"),
    (BRONZE_C,  "String Cleaning",
     "clean_string() — F.trim(F.regexp_replace(col, r'\\s+', ' ')) on all string columns"),
    (AGENT_C,   "Null Key Filtering",
     "Rows with NULL in primary-key columns (id, key, number) are dropped"),
    (SILVER_C,  "Critical Metric Nulls",
     "Rows with NULL in business-critical numeric columns (amounts, quantities) dropped"),
    (MUTED,     "AI-Generated Transforms",
     "New bronze tables with no silver counterpart → Llama 3.3 generates smart @dp.materialized_view"),
]
for i, (c, title, desc) in enumerate(rules):
    y = 2.1 + i * 0.66
    box(s4, 0.45, y, 0.08, 0.5, fill=c).line.fill.background()
    txt(s4, title, 0.7, y, 7.5, 0.28, size=12, bold=True, color=c)
    txt(s4, desc, 0.7, y + 0.3, 7.6, 0.28, size=10, color=MUTED)

# Right column — before/after
box(s4, 8.9, 1.5, 4.1, 4.9, fill=SURFACE)
txt(s4, "Data Quality Impact", 9.1, 1.6, 3.7, 0.4, size=15, bold=True, color=SILVER_C)
metrics = [
    ("dim_customers",  "18,484 → 18,469", "15 nulls removed"),
    ("fact_sales",     "60,398 → 60,379", "19 bad dates"),
    ("dim_products",   "295 → 295",       "no issues"),
    ("cust_info",      "filtered nulls",  "gender cleaned"),
]
for i, (tbl, change, note) in enumerate(metrics):
    y = 2.25 + i * 0.82
    box(s4, 9.0, y, 3.8, 0.7, fill=BG)
    txt(s4, tbl, 9.1, y+0.05, 3.6, 0.25, size=11, bold=True, color=WHITE)
    txt(s4, change, 9.1, y+0.28, 3.6, 0.22, size=11, color=GOLD_C)
    txt(s4, note, 9.1, y+0.48, 3.6, 0.2, size=10, color=MUTED)

txt(s4, "17 silver tables  ·  All views named silver_v2_<tablename>  ·  DLT incremental refresh",
    0.4, 6.6, 12.5, 0.4, size=11, color=SILVER_C, italic=True)

# ── Slide 5 — AI Agent ────────────────────────────────────────────────────────
s5 = prs.slides.add_slide(blank)
bg(s5)
box(s5, 0, 0, 0.08, 7.5, fill=AGENT_C).line.fill.background()
txt(s5, "AI Agent — Smart Orchestrator", 0.4, 0.3, 12, 0.6, size=28, bold=True, color=WHITE)
txt(s5, "Model: databricks-meta-llama-3-3-70b-instruct  ·  Databricks Foundation Model API  ·  No external API key",
    0.4, 0.9, 12.5, 0.35, size=12, color=MUTED)
divider(s5, 1.32)

steps_agent = [
    (1, "Silver Gap Detection",    BRONZE_C,
     "Finds bronze tables with no silver counterpart → calls Llama to write\n@dp.materialized_view cleaning code → appends to silver notebook → triggers silver DLT refresh"),
    (2, "Gold Gap Detection",      GOLD_C,
     "Checks catalog + SQL file for uncovered silver tables → Llama generates 2-4 gold SQL views\nper table with quality validation (semicolons, no lateral alias, unique names)"),
    ("2b", "Cross-Table Joins",   ACCENT,
     "Asks Llama if new tables share a key with existing gold views → suggests at most 2\ncross-table analytical views (e.g. sales joined to customer demographics)"),
    (3, "Selective Gold Refresh", AGENT_C,
     "Uses full_refresh_selection API to refresh ONLY the newly added views (2-5)\ninstead of all 75 views → reduces gold refresh from ~10 min to ~60 sec"),
    (4, "Dashboard Update",       SILVER_C,
     "Smart-updates 5 Lakeview dashboards (budget: 8 widgets/domain, 30 total)\nOnly adds widgets for new views — never overwrites existing ones"),
    (5, "GitHub Sync",            MUTED,
     "Exports updated silver notebook + gold SQL from Databricks workspace\nCommits them to GitHub via API so local repo stays in sync"),
]
for i, (num, title, c, desc) in enumerate(steps_agent):
    col = i % 2; row = i // 2
    bx = 0.3 + col * 6.5
    by = 1.55 + row * 1.78
    box(s5, bx, by, 6.2, 1.65, fill=SURFACE)
    box(s5, bx, by, 0.5, 0.4, fill=c).line.fill.background()
    txt(s5, str(num), bx+0.05, by+0.03, 0.4, 0.35, size=14, bold=True,
        color=BG, align=PP_ALIGN.CENTER)
    txt(s5, title, bx+0.6, by+0.05, 5.4, 0.35, size=13, bold=True, color=c)
    txt(s5, desc, bx+0.1, by+0.5, 5.9, 0.95, size=10, color=MUTED)

txt(s5, "Early exit: if no new tables detected → skips all steps and exits in <5 seconds",
    0.4, 7.0, 12.5, 0.35, size=11, color=GOLD_C, italic=True)

# ── Slide 6 — Gold Layer ──────────────────────────────────────────────────────
s6 = prs.slides.add_slide(blank)
bg(s6)
box(s6, 0, 0, 0.08, 7.5, fill=GOLD_C).line.fill.background()
txt(s6, "Gold Analytics Layer", 0.4, 0.3, 12, 0.6, size=28, bold=True, color=WHITE)
txt(s6, "75 Materialized Views  ·  4 Business Domains  ·  dataingestionproject.gold",
    0.4, 0.9, 12.5, 0.35, size=12, color=MUTED)
divider(s6, 1.32)

domains = [
    ("Customer Analytics", AGENT_C,
     ["gold_customer_summary", "gold_customer_demographics", "gold_customer_lifetime_value",
      "gold_customers_by_country", "gold_cust_info_*", "gold_loc_a101_by_country"]),
    ("Product Analytics", BRONZE_C,
     ["gold_product_summary", "gold_products_by_category", "gold_products_by_subcategory",
      "gold_prd_info_summary", "gold_px_cat_*", "gold_prd_info_by_line"]),
    ("Sales Analytics", GOLD_C,
     ["gold_sales_summary", "gold_monthly_sales", "gold_sales_by_country",
      "gold_sales_by_category", "gold_order_fulfillment", "gold_top_customers/products"]),
    ("Survey Analytics", ACCENT,
     ["gold_enterprise_survey_*", "gold_business_finance_*", "gold_climate_*",
      "gold_survey_*", "gold_student_performance_*", "gold_automobile_*"]),
]
for i, (domain, c, views) in enumerate(domains):
    col = i % 2; row = i // 2
    bx = 0.3 + col * 6.5
    by = 1.55 + row * 2.55
    box(s6, bx, by, 6.2, 2.3, fill=SURFACE)
    box(s6, bx, by, 6.2, 0.4, fill=c).line.fill.background()
    txt(s6, domain, bx+0.15, by+0.06, 5.8, 0.32, size=13, bold=True, color=BG)
    for j, v in enumerate(views):
        txt(s6, f"• {v}", bx+0.2, by+0.55 + j*0.28, 5.7, 0.26, size=10, color=MUTED)

txt(s6, "Selective refresh: only newly added views are refreshed per run  ·  UC quota: 490-table guard",
    0.4, 7.05, 12.5, 0.35, size=11, color=GOLD_C, italic=True)

# ── Slide 7 — Performance & Architecture ──────────────────────────────────────
s7 = prs.slides.add_slide(blank)
bg(s7)
txt(s7, "Performance & Architecture Decisions", 0.4, 0.3, 12, 0.6,
    size=28, bold=True, color=WHITE)
divider(s7, 1.0)

# Timing comparison
box(s7, 0.3, 1.2, 5.8, 5.2, fill=SURFACE)
txt(s7, "Job Timing — Before vs After Optimization", 0.55, 1.3, 5.3, 0.4,
    size=14, bold=True, color=ACCENT)

timing = [
    ("Bronze",    "~120s",  "~72s",   BRONZE_C),
    ("Silver",    "~172s",  "~172s",  SILVER_C),
    ("AI Agent",  "~100s",  "~348s",  AGENT_C),
    ("Gold (DLT)","~675s",  "inside agent", GOLD_C),
    ("Total",     "~18 min","~10 min", WHITE),
]
header_y = 1.85
txt(s7, "Task", 0.55, header_y, 2.0, 0.28, size=11, bold=True, color=MUTED)
txt(s7, "Before", 2.55, header_y, 1.5, 0.28, size=11, bold=True, color=MUTED)
txt(s7, "After", 4.05, header_y, 1.8, 0.28, size=11, bold=True, color=MUTED)
for i, (task, before, after, c) in enumerate(timing):
    y = 2.25 + i * 0.72
    if task == "Total":
        box(s7, 0.35, y-0.05, 5.7, 0.75, fill=BG).line.fill.background()
    txt(s7, task, 0.55, y, 1.9, 0.35, size=12, bold=(task=="Total"), color=c)
    txt(s7, before, 2.55, y, 1.4, 0.35, size=12, color=MUTED)
    txt(s7, after, 4.05, y, 1.8, 0.35, size=12, bold=(task=="Total"),
        color=GOLD_C if task=="Total" else WHITE)

# Key decisions
box(s7, 6.5, 1.2, 6.5, 5.2, fill=SURFACE)
txt(s7, "Key Architecture Decisions", 6.75, 1.3, 6.0, 0.4,
    size=14, bold=True, color=ACCENT)
decisions = [
    (GOLD_C,   "Selective gold refresh",
     "full_refresh_selection refreshes 2-5 new views instead of all 75 → saves ~600s per run"),
    (AGENT_C,  "Agent controls gold",
     "Gold DLT task removed from job; agent triggers selectively and waits for completion"),
    (BRONZE_C, "Incremental bronze",
     "changed_files param filters to only new CSVs; full re-ingestion only on manual runs"),
    (SILVER_C, "Early exit",
     "Agent exits in <5s if no new tables detected — zero wasted compute"),
    (MUTED,    "Metastore quota guard",
     "Both bronze and agent check system.information_schema.tables before creating objects"),
    (ACCENT,   "GitHub auto-sync",
     "Agent commits updated silver/gold files to GitHub after each run via GitHub API"),
    (GOLD_C,   "SQL quality validation",
     "Post-generation checks: semicolons, no duplicate names, no lateral alias in WINDOW"),
]
for i, (c, title, desc) in enumerate(decisions):
    y = 1.85 + i * 0.65
    box(s7, 6.6, y, 0.07, 0.45, fill=c).line.fill.background()
    txt(s7, title, 6.8, y, 5.9, 0.25, size=11, bold=True, color=c)
    txt(s7, desc, 6.8, y+0.28, 5.9, 0.28, size=10, color=MUTED)

# ── Save ──────────────────────────────────────────────────────────────────────
output = r'c:\Users\uik11839\Documents\databricksProjectDataSource\DataProcessing_Pipeline_Architecture.pptx'
prs.save(output)
print(f'Saved: {output}')
